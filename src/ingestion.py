"""
ingestion.py
Responsibility: Asmit

This module handles the extraction and cleaning of text from raw PDF and PPTX documents.
"""
import os
import tempfile
from typing import List
from io import BytesIO

from pypdf import PdfReader
from pptx import Presentation
from langchain_core.documents import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter

def extract_text_from_pdf(file_stream) -> str:
    """Extracts text from a PDF file stream."""
    text = ""
    try:
        pdf_reader = PdfReader(file_stream)
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        print(f"Error reading PDF: {e}")
    return text

def extract_text_from_pptx(file_stream) -> str:
    """Extracts text from a PPTX file stream."""
    text = ""
    try:
        prs = Presentation(file_stream)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX: {e}")
    return text

def load_documents(uploaded_files) -> List[Document]:
    """
    Processes a list of uploaded files (PDF or PPTX) and returns a list of split Document objects.
    """
    raw_text = ""
    
    for uploaded_file in uploaded_files:
        # Streamlit UploadedFile object has .name and .read() (or is a stream)
        # We can pass the UploadedFile directly to the readers in many cases, 
        # but safely we might need to reset position or handle bytes.
        
        file_name = uploaded_file.name.lower()
        
        # We need to ensure we are at the start of the file
        uploaded_file.seek(0)
        
        if file_name.endswith('.pdf'):
            text = extract_text_from_pdf(uploaded_file)
            raw_text += text + "\n\n"
        elif file_name.endswith('.pptx') or file_name.endswith('.ppt'):
            text = extract_text_from_pptx(uploaded_file)
            raw_text += text + "\n\n"
        else:
            # Fallback for text files if needed, or skip
            pass

    # Split the text
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        length_function=len,
        is_separator_regex=False,
    )
    
    texts = text_splitter.split_text(raw_text)
    
    # Create Document objects
    docs = [Document(page_content=t) for t in texts]
    
    return docs

if __name__ == "__main__":
    # Test logic (Mocking a file update if this script is run directly is hard without true files)
    print("Run this module from app.py by passing UploadedFiles.")
