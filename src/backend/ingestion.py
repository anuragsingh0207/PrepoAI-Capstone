from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from pypdf import PdfReader
from pptx import Presentation
import os


def load_documents(uploaded_files, chunk_size=500, overlap=50):
    """Process each uploaded file independently and attach metadata."""
    documents = []
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        length_function=len,
        is_separator_regex=False,
    )

    for uploaded_file in uploaded_files:
        name = uploaded_file.name
        extension = os.path.splitext(name)[1].lower()
        uploaded_file.seek(0)

        # --- PDF ---
        if extension == ".pdf":
            pdf_reader = PdfReader(uploaded_file)
            for i, page in enumerate(pdf_reader.pages):
                text = page.extract_text() or ""
                for chunk in splitter.split_text(text):
                    documents.append(
                        Document(
                            page_content=chunk.strip(),
                            metadata={"source": name, "page": i + 1, "type": "pdf"},
                        )
                    )

        # --- PPTX ---
        elif extension in [".pptx", ".ppt"]:
            prs = Presentation(uploaded_file)
            for i, slide in enumerate(prs.slides):
                slide_text = "\n".join(
                    [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                )
                for chunk in splitter.split_text(slide_text):
                    documents.append(
                        Document(
                            page_content=chunk.strip(),
                            metadata={"source": name, "slide": i + 1, "type": "ppt"},
                        )
                    )

        else:
            print(f"Unsupported file type: {name}")

    print(f"✅ Prepared {len(documents)} document chunks with metadata.")
    return documents

